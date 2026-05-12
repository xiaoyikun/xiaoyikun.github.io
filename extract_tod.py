"""按页码分组提取 TOD 4 个场景的大图到 portfolio/"""
import os
from pptx import Presentation

PPTX = r"C:\Users\pixxiao\Desktop\肖奕堃广东工业大学作品集.pptx"
OUT = r"C:\Users\pixxiao\Desktop\升级版作品集\src\assets\portfolio"

# (页码范围, 前缀)
GROUPS = [
    (range(5, 9),   "tod-s1"),  # 5-8
    (range(9, 13),  "tod-s2"),  # 9-12
    (range(13, 15), "tod-s3"),  # 13-14
    (range(15, 16), "tod-s4"),  # 15 (只取第1张)
]

# 先清掉旧的 tod-*.png
for f in os.listdir(OUT):
    if f.startswith("tod-") or f == "tod-01.png":
        if f.startswith("tod-") and not f.startswith("tod-s"):
            os.remove(os.path.join(OUT, f))
            print("removed", f)

p = Presentation(PPTX)
SW, SH = p.slide_width, p.slide_height
MIN = 0.08

for rng, prefix in GROUPS:
    n = 0
    take_one = (prefix == "tod-s4")
    for idx in rng:
        sl = p.slides[idx-1]
        pics = []
        for sh in sl.shapes:
            if sh.shape_type != 13:
                continue
            area = (sh.width or 0)*(sh.height or 0)/(SW*SH)
            if area < MIN:
                continue
            pics.append((sh.top or 0, sh.left or 0, sh))
        pics.sort(key=lambda x: (x[0], x[1]))
        for _, _, sh in pics:
            n += 1
            out = os.path.join(OUT, f"{prefix}-{n:02d}.png")
            with open(out, "wb") as fp:
                fp.write(sh.image.blob)
            print(f"slide {idx} -> {prefix}-{n:02d}.png")
            if take_one:
                break
        if take_one and n >= 1:
            break
