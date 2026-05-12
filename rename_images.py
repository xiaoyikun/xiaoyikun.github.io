"""
按 PPTX 幻灯片顺序提取大图（跳过小装饰图），按作品分组重命名到 src/assets/portfolio/
"""
import os
import shutil
from pptx import Presentation
from pptx.util import Emu

PPTX = r"C:\Users\pixxiao\Desktop\肖奕堃广东工业大学作品集.pptx"
OUT = r"C:\Users\pixxiao\Desktop\升级版作品集\src\assets\portfolio"

# 幻灯片范围 -> 作品前缀
RANGES = [
    (range(4, 16), "tod"),         # 4-15: TOD
    (range(16, 24), "relight"),    # 16-23: 同场景重打光
    (range(24, 31), "jp-garden"),  # 24-30: 日式庭院 PBR
    (range(31, 33), "char-light"), # 31-32: 角色灯光
    (range(34, 35), "bamboo"),     # 34: 竹林
    (range(35, 37), "env"),        # 35-36: 环境美术速写
    (range(37, 39), "church"),     # 37-38: 教堂
]

def slide_prefix(idx):
    for r, name in RANGES:
        if idx in r:
            return name
    return None

# 清空旧目录
if os.path.exists(OUT):
    shutil.rmtree(OUT)
os.makedirs(OUT, exist_ok=True)

p = Presentation(PPTX)

# 幻灯片宽高（EMU）
SW = p.slide_width
SH = p.slide_height
# 只要面积 >= 整页 8% 的图（过滤角标、logo、装饰条）
MIN_AREA_RATIO = 0.08

counters = {}  # prefix -> next index
seen_hashes = set()  # 去重相同的图

for idx, sl in enumerate(p.slides, start=1):
    prefix = slide_prefix(idx)
    if prefix is None:
        continue

    # 收集本页所有图片，按从上到下、从左到右排序
    pics = []
    for shape in sl.shapes:
        if shape.shape_type != 13:  # PICTURE
            continue
        try:
            w = shape.width or 0
            h = shape.height or 0
        except Exception:
            w = h = 0
        area_ratio = (w * h) / (SW * SH) if SW and SH else 0
        if area_ratio < MIN_AREA_RATIO:
            continue
        top = shape.top or 0
        left = shape.left or 0
        pics.append((top, left, shape))

    pics.sort(key=lambda x: (x[0], x[1]))

    for _, _, shape in pics:
        blob = shape.image.blob
        h = hash(blob)
        if h in seen_hashes:
            continue
        seen_hashes.add(h)

        ext = shape.image.ext or "png"
        counters[prefix] = counters.get(prefix, 0) + 1
        n = counters[prefix]
        out_name = f"{prefix}-{n:02d}.{ext}"
        out_path = os.path.join(OUT, out_name)
        with open(out_path, "wb") as f:
            f.write(blob)
        print(f"slide {idx:>2} -> {out_name}  ({len(blob)//1024} KB)")

print("\n==== summary ====")
for k, v in counters.items():
    print(f"{k}: {v} 张")
